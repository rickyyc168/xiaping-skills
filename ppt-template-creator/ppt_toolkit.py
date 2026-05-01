#!/usr/bin/env python3
"""
PPT Template Creator v3 — Design System + Toolkit.

Commands:
    parse    - Parse existing PPT to JSON (extract all styles)
    generate - Generate PPT from JSON spec (auto-adds page badges)
    inject   - Template + content JSON → new PPT
    qa       - Verify PPT: extract text, check for issues
    demo     - Generate demo with all 5 page types

Usage:
    python3 ppt_toolkit.py parse   input.pptx -o styles.json
    python3 ppt_toolkit.py generate spec.json  -o output.pptx
    python3 ppt_toolkit.py inject  template.pptx content.json -o result.pptx
    python3 ppt_toolkit.py qa      output.pptx
    python3 ppt_toolkit.py demo    -o demo.pptx
"""

import json
import sys
import argparse
import re
import zipfile
import io
from pathlib import Path
from copy import deepcopy

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
    from pptx.enum.dml import MSO_FILL, MSO_LINE_DASH_STYLE
    try:
        from pptx.enum.dml import MSO_THEME_COLOR
    except ImportError:
        MSO_THEME_COLOR = None
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


# ═══════════════════════════════════════════════════════════════════════════════
#  COMMON HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def hex_to_rgb(hex_str):
    if not hex_str:
        return None
    h = str(hex_str).lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def rgb_to_hex(color_obj):
    if color_obj is None:
        return None
    if isinstance(color_obj, (tuple, list)) and len(color_obj) == 3:
        return f"{color_obj[0]:02X}{color_obj[1]:02X}{color_obj[2]:02X}"
    try:
        r, g, b = color_obj
        return f"{r:02X}{g:02X}{b:02X}"
    except:
        return str(color_obj)


ALIGN_MAP = {
    "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY,
    "PP_ALIGN.LEFT": PP_ALIGN.LEFT, "PP_ALIGN.CENTER": PP_ALIGN.CENTER,
    "PP_ALIGN.RIGHT": PP_ALIGN.RIGHT, "PP_ALIGN.JUSTIFY": PP_ALIGN.JUSTIFY,
}


def get_alignment(val):
    if val is None:
        return PP_ALIGN.LEFT
    return ALIGN_MAP.get(str(val), PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════════
#  PAGE BADGE
# ═══════════════════════════════════════════════════════════════════════════════

def add_page_badge(slide, pres, page_num, theme, style="circle"):
    """Add a page number badge to bottom-right corner."""
    accent = hex_to_rgb(theme.get("accent", "E94560"))
    white = RGBColor(0xFF, 0xFF, 0xFF)

    if style == "pill":
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(9.1), Inches(5.15), Inches(0.6), Inches(0.35)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = accent
        shape.line.fill.background()
        try:
            shape.adjustments[0] = 0.5
        except:
            pass
        txBox = slide.shapes.add_textbox(
            Inches(9.1), Inches(5.15), Inches(0.6), Inches(0.35)
        )
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(9.3), Inches(5.1), Inches(0.4), Inches(0.4)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = accent
        shape.line.fill.background()
        txBox = slide.shapes.add_textbox(
            Inches(9.3), Inches(5.1), Inches(0.4), Inches(0.4)
        )

    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(page_num)
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.CENTER


# ═══════════════════════════════════════════════════════════════════════════════
#  PARSE — PPT → JSON
# ═══════════════════════════════════════════════════════════════════════════════

def parse_fill(fill):
    try:
        return {
            "type": str(MSO_FILL(fill.type)) if fill.type else None,
            "color": rgb_to_hex(fill.fore_color.rgb) if fill.fore_color and fill.fore_color.rgb else None,
        }
    except:
        return {"type": None, "color": None}


def parse_line(line):
    try:
        return {
            "color": rgb_to_hex(line.color.rgb) if line.color and line.color.rgb else None,
            "width": int(line.width) if line.width else 0,
            "dash_style": str(MSO_LINE_DASH_STYLE(line.dash_style)) if line.dash_style else None,
        }
    except:
        return {"color": None, "width": 0, "dash_style": None}


def parse_text_style(text_frame):
    paragraphs = []
    for para in text_frame.paragraphs:
        runs = []
        for run in para.runs:
            color_info = {"type": "rgb", "hex": None}
            try:
                if run.font.color and run.font.color.rgb:
                    color_info = {"type": "rgb", "hex": rgb_to_hex(run.font.color.rgb)}
                elif MSO_THEME_COLOR and run.font.color and run.font.color.type == MSO_THEME_COLOR:
                    color_info = {"type": "theme", "theme_color": str(run.font.color.theme_color)}
            except:
                pass
            runs.append({
                "text": run.text,
                "font": {
                    "name": run.font.name,
                    "size": run.font.size.pt if run.font.size else None,
                    "bold": run.font.bold,
                    "italic": run.font.italic,
                    "color": color_info,
                }
            })
        paragraphs.append({
            "text": para.text,
            "alignment": str(para.alignment) if para.alignment else None,
            "level": para.level,
            "runs": runs,
        })
    return {"paragraphs": paragraphs}


def parse_shape(shape):
    data = {
        "type": int(shape.shape_type) if shape.shape_type else None,
        "type_name": str(MSO_SHAPE_TYPE(shape.shape_type)) if shape.shape_type else "UNKNOWN",
        "name": shape.name,
        "left": int(shape.left), "top": int(shape.top),
        "width": int(shape.width), "height": int(shape.height),
        "rotation": shape.rotation,
    }
    try:
        data["fill"] = parse_fill(shape.fill)
    except:
        data["fill"] = {"type": None, "color": None}
    try:
        data["line"] = parse_line(shape.line)
    except:
        data["line"] = {"color": None, "width": 0, "dash_style": None}
    if shape.has_text_frame:
        data["text"] = shape.text_frame.text
        data["text_style"] = parse_text_style(shape.text_frame)
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        table = shape.table
        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cells.append({
                    "text": cell.text.strip(),
                    "row_span": cell.span_height if hasattr(cell, 'span_height') else 1,
                    "col_span": cell.span_width if hasattr(cell, 'span_width') else 1,
                    "is_merge_origin": cell.is_merge_origin if hasattr(cell, 'is_merge_origin') else False,
                })
            rows.append(cells)
        data["table"] = rows
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        data["image"] = {"width": int(shape.width), "height": int(shape.height)}
    return data


def parse_presentation(pptx_path):
    prs = Presentation(pptx_path)
    result = {
        "source": str(pptx_path),
        "width": int(prs.slide_width), "height": int(prs.slide_height),
        "width_inches": round(prs.slide_width / 914400, 3),
        "height_inches": round(prs.slide_height / 914400, 3),
        "slides": [],
    }
    for slide_idx, slide in enumerate(prs.slides):
        slide_data = {"slide_number": slide_idx + 1, "slide_id": slide.slide_id, "shapes": []}
        for shape in slide.shapes:
            slide_data["shapes"].append(parse_shape(shape))
        result["slides"].append(slide_data)
    print(f"Parsed: {pptx_path}")
    print(f"  Size: {result['width_inches']}\" × {result['height_inches']}\"")
    print(f"  Slides: {len(result['slides'])}")
    for s in result["slides"]:
        print(f"    Slide {s['slide_number']}: {len(s['shapes'])} shapes")
    return result


# ═══════════════════════════════════════════════════════════════════════════════
#  GENERATE — JSON Spec → PPT
# ═══════════════════════════════════════════════════════════════════════════════

def _apply_text(tf, spec, text="", default_color="333333"):
    """Apply text to a text frame with formatting."""
    font_size = Pt(spec.get("font_size", 18))
    bold = spec.get("bold", False)
    italic = spec.get("italic", False)
    align = get_alignment(spec.get("align", "left"))
    color = hex_to_rgb(spec.get("color", default_color))
    font_name = spec.get("font_name")

    if text:
        lines = str(text).split("\n")
        tf.paragraphs[0].text = lines[0]
        tf.paragraphs[0].font.size = font_size
        tf.paragraphs[0].font.bold = bold
        tf.paragraphs[0].font.italic = italic
        tf.paragraphs[0].font.color.rgb = color
        tf.paragraphs[0].alignment = align
        if font_name:
            tf.paragraphs[0].font.name = font_name
        for line in lines[1:]:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = font_size
            p.font.bold = bold
            p.font.italic = italic
            p.font.color.rgb = color
            p.alignment = align
            if font_name:
                p.font.name = font_name
    else:
        tf.paragraphs[0].font.size = font_size
        tf.paragraphs[0].font.color.rgb = color
    return tf


def add_textbox(slide, spec, text="", default_color="333333"):
    txBox = slide.shapes.add_textbox(
        Inches(spec.get("left", 0.5)), Inches(spec.get("top", 0.5)),
        Inches(spec.get("width", 9)), Inches(spec.get("height", 1))
    )
    txBox.text_frame.word_wrap = spec.get("word_wrap", True)
    return _apply_text(txBox.text_frame, spec, text, default_color)


def add_bullet_list(slide, spec, items, default_color="333333"):
    txBox = slide.shapes.add_textbox(
        Inches(spec.get("left", 0.5)), Inches(spec.get("top", 0.5)),
        Inches(spec.get("width", 9)), Inches(spec.get("height", 1))
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    font_size = Pt(spec.get("font_size", 16))
    color = hex_to_rgb(spec.get("color", default_color))
    bold = spec.get("bold", False)
    font_name = spec.get("font_name")
    bullet = spec.get("bullet", "•")

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, dict):
            p.text = f"{bullet} {item.get('heading', '')}"
            p.font.size = font_size
            p.font.bold = True
            p.font.color.rgb = color
            for sub in item.get("items", []):
                sp = tf.add_paragraph()
                sp.text = f"  ◦ {sub}"
                sp.font.size = Pt(spec.get("font_size", 16) - 2)
                sp.font.color.rgb = color
                if font_name:
                    sp.font.name = font_name
        else:
            p.text = f"{bullet} {item}"
            p.font.size = font_size
            p.font.bold = bold
            p.font.color.rgb = color
        if font_name:
            p.font.name = font_name
    return tf


def add_table_shape(slide, spec, data):
    left = Inches(spec.get("left", 0.5))
    top = Inches(spec.get("top", 1.5))
    width = Inches(spec.get("width", 12))
    height = Inches(spec.get("height", 4))
    rows = len(data)
    cols = max(len(r) for r in data) if data else 1
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    hdr_fs = Pt(spec.get("header_font_size", 14))
    cell_fs = Pt(spec.get("cell_font_size", 12))
    hdr_color = hex_to_rgb(spec.get("header_color", "FFFFFF"))
    hdr_bg = hex_to_rgb(spec.get("header_bg", "1A1A2E"))
    cell_color = hex_to_rgb(spec.get("cell_color", "333333"))

    for r_idx, row in enumerate(data):
        for c_idx, cell_data in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            text = cell_data.get("text", str(cell_data)) if isinstance(cell_data, dict) else str(cell_data)
            cell.text = text
            for p in cell.text_frame.paragraphs:
                if r_idx == 0:
                    p.font.size = hdr_fs
                    p.font.bold = True
                    p.font.color.rgb = hdr_color
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = hdr_bg
                else:
                    p.font.size = cell_fs
                    p.font.color.rgb = cell_color
    return table_shape


def add_rect_shape(slide, spec):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(spec.get("left", 0)), Inches(spec.get("top", 0)),
        Inches(spec.get("width", 1)), Inches(spec.get("height", 1))
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(spec.get("fill", "CCCCCC"))
    shape.line.fill.background()
    if spec.get("rounded"):
        try:
            shape.adjustments[0] = 0.05
        except:
            pass
    return shape


def set_slide_bg(slide, color_str):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_str)


ELEMENT_BUILDERS = {
    "text": lambda s, e, sd, tc: add_textbox(s, e, e.get("text", sd.get(e.get("content_key", ""), "")), tc),
    "title": lambda s, e, sd, tc: add_textbox(s, e, str(sd.get(e.get("content_key", "title"), e.get("default", ""))), tc),
    "subtitle": lambda s, e, sd, tc: add_textbox(s, e, str(sd.get(e.get("content_key", "subtitle"), e.get("default", ""))), tc),
    "body": lambda s, e, sd, tc: (
        add_bullet_list(s, e, sd.get(e.get("content_key", "body"), []), tc)
        if isinstance(sd.get(e.get("content_key", "body"), e.get("default", "")), list)
        else add_textbox(s, e, str(sd.get(e.get("content_key", "body"), e.get("default", ""))), tc)
    ),
    "rect": lambda s, e, sd, tc: add_rect_shape(s, e),
    "divider": lambda s, e, sd, tc: add_rect_shape(s, {**e, "height": e.get("height", 0.04)}),
    "image": lambda s, e, sd, tc: (
        s.shapes.add_picture(sd.get(e.get("content_key", "image")),
                             Inches(e["left"]), Inches(e["top"]),
                             Inches(e["width"]), Inches(e["height"]))
        if Path(sd.get(e.get("content_key", "image"), "")).exists() else None
    ),
    "table": lambda s, e, sd, tc: add_table_shape(s, e, sd.get(e.get("content_key", "table"), [])),
}


def build_slide_from_spec(prs, layout_spec, slide_data, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    default_tc = theme.get("text", theme.get("secondary", "333333"))
    text_color = slide_data.get("text_color", layout_spec.get("text_color", default_tc))
    bg = slide_data.get("background", layout_spec.get("background"))
    if bg:
        set_slide_bg(slide, bg)
    for elem in layout_spec.get("elements", []):
        etype = elem.get("type", "text")
        builder = ELEMENT_BUILDERS.get(etype)
        if builder:
            builder(slide, elem, slide_data, text_color)
    return slide


def generate_from_spec(spec, add_badges=True):
    prs = Presentation()
    prs.slide_width = Inches(spec.get("width_inches", 13.333))
    prs.slide_height = Inches(spec.get("height_inches", 7.5))

    theme = spec.get("theme", {})
    # Normalize theme: support both "text" and "secondary" for text color
    if "text" not in theme and "secondary" in theme:
        theme["text"] = theme["secondary"]

    layouts = spec.get("layouts", [])
    slides = spec.get("slides", [])

    print(f"Generating: {spec.get('name', 'Untitled')}")
    print(f"  Size: {spec.get('width_inches', 13.333)}\" × {spec.get('height_inches', 7.5)}\"")
    print(f"  Theme: {json.dumps({k: v for k, v in theme.items() if k in ['primary','secondary','accent','light','bg']})}")
    print(f"  Layouts: {len(layouts)}, Slides: {len(slides)}")

    page_num = 0
    for i, sd in enumerate(slides):
        li = sd.get("layout", 0)
        if li >= len(layouts):
            print(f"  ⚠ Slide {i+1}: layout {li} out of range, using 0")
            li = 0

        layout_spec = layouts[li]
        layout_type = layout_spec.get("type", "content")

        build_slide_from_spec(prs, layout_spec, sd, theme)

        # Add page badge (skip cover pages)
        if add_badges and layout_type != "cover":
            page_num += 1
            add_page_badge(prs.slides[-1], prs, page_num, theme, style=spec.get("badge_style", "circle"))

        type_label = f"[{layout_type}]" if layout_type != "content" else ""
        print(f"  Slide {i+1}: layout={li} ({layout_spec.get('name', '?')}) {type_label}")

    return prs


# ═══════════════════════════════════════════════════════════════════════════════
#  INJECT — Template + Content → PPT
# ═══════════════════════════════════════════════════════════════════════════════

def _replace_text_preserving_runs(text_frame, new_text):
    if not text_frame.paragraphs:
        return
    if isinstance(new_text, list):
        text_frame.paragraphs[0].text = str(new_text[0]) if new_text else ""
        src_run = text_frame.paragraphs[0].runs[0] if text_frame.paragraphs[0].runs else None
        for item in new_text[1:]:
            p = text_frame.add_paragraph()
            p.text = str(item)
            if src_run:
                p.font.name = src_run.font.name
                p.font.size = src_run.font.size
                p.font.bold = src_run.font.bold
                p.font.italic = src_run.font.italic
                try:
                    if src_run.font.color and src_run.font.color.rgb:
                        p.font.color.rgb = src_run.font.color.rgb
                except:
                    pass
    else:
        if text_frame.paragraphs[0].runs:
            text_frame.paragraphs[0].runs[0].text = str(new_text)
        else:
            text_frame.paragraphs[0].text = str(new_text)


def inject_content(template_path, content_path, output_path):
    with open(content_path, "r", encoding="utf-8") as f:
        content = json.load(f)
    prs = Presentation(template_path)
    slides_content = content.get("slides", [])

    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx >= len(slides_content):
            break
        sc = slides_content[slide_idx]

        if "shapes" in sc:
            name_map = {s.get("name", ""): s for s in sc["shapes"]}
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                match = name_map.get(shape.name)
                if match and "text" in match:
                    _replace_text_preserving_runs(shape.text_frame, match["text"])

        if "texts" in sc:
            text_idx = 0
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if text_idx < len(sc["texts"]):
                    _replace_text_preserving_runs(shape.text_frame, sc["texts"][text_idx])
                    text_idx += 1

    prs.save(output_path)
    print(f"Injected: {template_path} + {content_path} → {output_path}")
    return prs


# ═══════════════════════════════════════════════════════════════════════════════
#  QA — Verify PPTX
# ═══════════════════════════════════════════════════════════════════════════════

PLACEHOLDER_PATTERNS = [
    r"lorem\s*ipsum", r"placeholder", r"xxxx", r"click to add",
    r"this\s*(page|slide)\s*(is\s*)?a\s*layout", r"\[your\s+text\s+here\]",
    r"sample\s+text", r"edit\s+this", r"type\s+here",
]

def qa_presentation(pptx_path):
    """Run QA checks on a PPTX file. Returns list of issues."""
    issues = []
    prs = Presentation(pptx_path)

    print(f"QA: {pptx_path}")
    print(f"  Slides: {len(prs.slides)}")
    print(f"  Size: {round(prs.slide_width/914400,3)}\" × {round(prs.slide_height/914400,3)}\"")
    print()

    all_text = []
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        shape_count = 0
        text_shapes = 0
        has_visual = False

        for shape in slide.shapes:
            shape_count += 1
            if shape.has_text_frame:
                text_shapes += 1
                text = shape.text_frame.text.strip()
                all_text.append((slide_num, shape.name, text))

                # Check for placeholder text
                for pattern in PLACEHOLDER_PATTERNS:
                    if re.search(pattern, text, re.IGNORECASE):
                        issues.append(f"Slide {slide_num}: Placeholder text found in '{shape.name}': '{text[:50]}...'")

                # Check for empty required text
                if not text and shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    pass  # Empty textboxes are ok for decorative purposes

            # Check for visual elements
            visual_types = {'PICTURE', 'TABLE', 'CHART', 'MEDIA', 'WEB_VIDEO'}
            shape_type_name = str(shape.shape_type).split('.')[-1] if hasattr(shape.shape_type, '__str__') else ''
            if shape_type_name in visual_types or shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_visual = True
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                has_visual = True

        # Content slides should have visual elements
        if slide_idx > 0 and shape_count > 0 and not has_visual and text_shapes > 0:
            # Only warn if there are more than 2 text shapes (likely a content slide)
            if text_shapes >= 2:
                issues.append(f"Slide {slide_num}: Content slide with no visual elements (only text shapes)")

    # Print text extraction
    print("─── Text Extraction ───")
    for slide_num, shape_name, text in all_text:
        if text:
            preview = text[:80] + ("..." if len(text) > 80 else "")
            print(f"  [{slide_num}] {shape_name}: {preview}")
    print()

    # Print issues
    if issues:
        print(f"─── Issues Found: {len(issues)} ───")
        for issue in issues:
            print(f"  ⚠ {issue}")
    else:
        print("─── No issues found ───")
        print("  (But look harder — first pass rarely finds everything)")

    return issues


# ═══════════════════════════════════════════════════════════════════════════════
#  DEMO — All 5 Page Types
# ═══════════════════════════════════════════════════════════════════════════════

DEMO_SPEC = {
    "name": "Design System Demo — 5 Page Types",
    "width_inches": 13.333,
    "height_inches": 7.5,
    "theme": {
        "primary": "1A1A2E", "secondary": "16213E", "accent": "E94560",
        "light": "F5F5F5", "bg": "FFFFFF"
    },
    "badge_style": "circle",
    "layouts": [
        # 0: Cover
        {
            "name": "Cover", "type": "cover",
            "background": "1A1A2E", "text_color": "FFFFFF",
            "elements": [
                {"type": "title", "left": 0.8, "top": 2.0, "width": 11.7, "height": 1.5,
                 "font_size": 44, "bold": True, "align": "center"},
                {"type": "subtitle", "left": 0.8, "top": 3.8, "width": 11.7, "height": 0.8,
                 "font_size": 20, "align": "center"},
                {"type": "divider", "left": 5.0, "top": 3.6, "width": 3.3, "height": 0.04, "fill": "E94560"},
            ]
        },
        # 1: Table of Contents
        {
            "name": "TOC", "type": "toc",
            "elements": [
                {"type": "title", "left": 0.5, "top": 0.3, "width": 12.3, "height": 0.9,
                 "font_size": 32, "bold": True, "color": "1A1A2E"},
                {"type": "divider", "left": 0.5, "top": 1.3, "width": 2.0, "height": 0.04, "fill": "E94560"},
                {"type": "body", "left": 0.5, "top": 1.8, "width": 12.3, "height": 5.0,
                 "font_size": 18, "color": "16213E", "bullet": "▸"},
                {"type": "rect", "left": 0.3, "top": 1.8, "width": 0.06, "height": 4.5, "fill": "E94560"},
            ]
        },
        # 2: Section Divider
        {
            "name": "Section Divider", "type": "section",
            "background": "E94560", "text_color": "FFFFFF",
            "elements": [
                {"type": "text", "left": 0.8, "top": 1.5, "width": 11.7, "height": 1.5,
                 "font_size": 72, "bold": True, "align": "center", "content_key": "section_num"},
                {"type": "title", "left": 0.8, "top": 3.0, "width": 11.7, "height": 1.0,
                 "font_size": 32, "bold": True, "align": "center"},
                {"type": "subtitle", "left": 1.5, "top": 4.2, "width": 10.3, "height": 0.6,
                 "font_size": 16, "align": "center", "italic": True},
            ]
        },
        # 3: Content — Bullet List
        {
            "name": "Content Bullets", "type": "content",
            "elements": [
                {"type": "title", "left": 0.5, "top": 0.3, "width": 12.3, "height": 0.9,
                 "font_size": 28, "bold": True, "color": "1A1A2E"},
                {"type": "divider", "left": 0.5, "top": 1.3, "width": 2.0, "height": 0.04, "fill": "E94560"},
                {"type": "body", "left": 0.5, "top": 1.6, "width": 12.3, "height": 5.5,
                 "font_size": 16, "color": "16213E"},
            ]
        },
        # 4: Content — Two Column
        {
            "name": "Two Column", "type": "content",
            "elements": [
                {"type": "title", "left": 0.5, "top": 0.3, "width": 12.3, "height": 0.9,
                 "font_size": 28, "bold": True, "color": "1A1A2E"},
                {"type": "divider", "left": 0.5, "top": 1.3, "width": 2.0, "height": 0.04, "fill": "E94560"},
                {"type": "rect", "left": 0.5, "top": 1.6, "width": 5.8, "height": 5.2, "fill": "F5F5F5"},
                {"type": "body", "left": 0.8, "top": 1.8, "width": 5.2, "height": 4.8,
                 "font_size": 14, "color": "16213E", "content_key": "left"},
                {"type": "rect", "left": 6.8, "top": 1.6, "width": 5.8, "height": 5.2, "fill": "F5F5F5"},
                {"type": "body", "left": 7.1, "top": 1.8, "width": 5.2, "height": 4.8,
                 "font_size": 14, "color": "16213E", "content_key": "right"},
            ]
        },
        # 5: Content — Table
        {
            "name": "Data Table", "type": "content",
            "elements": [
                {"type": "title", "left": 0.5, "top": 0.3, "width": 12.3, "height": 0.9,
                 "font_size": 28, "bold": True, "color": "1A1A2E"},
                {"type": "divider", "left": 0.5, "top": 1.3, "width": 2.0, "height": 0.04, "fill": "E94560"},
                {"type": "table", "left": 0.5, "top": 1.6, "width": 12.3, "height": 5.0,
                 "header_bg": "1A1A2E", "header_color": "FFFFFF"},
            ]
        },
        # 6: Summary / Closing
        {
            "name": "Closing", "type": "summary",
            "background": "1A1A2E", "text_color": "FFFFFF",
            "elements": [
                {"type": "title", "left": 0.8, "top": 1.5, "width": 11.7, "height": 1.2,
                 "font_size": 40, "bold": True, "align": "center"},
                {"type": "divider", "left": 5.0, "top": 2.9, "width": 3.3, "height": 0.04, "fill": "E94560"},
                {"type": "body", "left": 1.5, "top": 3.2, "width": 10.3, "height": 2.5,
                 "font_size": 18, "align": "center", "bullet": "✓"},
                {"type": "subtitle", "left": 0.8, "top": 5.8, "width": 11.7, "height": 0.6,
                 "font_size": 14, "align": "center"},
            ]
        },
    ],
    "slides": [
        # Type A: Cover
        {"layout": 0, "title": "年度业务总结报告", "subtitle": "2025 财年 | 产品技术部"},
        # Type B: Table of Contents
        {"layout": 1, "title": "目录",
         "body": ["01  项目背景与目标", "02  核心成果与数据", "03  团队与协作", "04  未来规划与展望"]},
        # Type C: Section Divider
        {"layout": 2, "section_num": "01", "title": "项目背景与目标", "subtitle": "从战略规划到落地执行的全景回顾"},
        # Type D: Content — Bullets
        {"layout": 3, "title": "项目背景",
         "body": ["市场调研覆盖 12 个行业，访谈 200+ 用户", "技术选型完成，确定微服务架构方案",
                  "核心团队 15 人组建到位", "Q1 完成 MVP，Q2 启动内测"]},
        # Type D: Content — Two Column
        {"layout": 4, "title": "方案对比",
         "left": ["方案 A：自建平台", "• 优势：完全可控\n• 成本：高\n• 周期：12 个月\n• 风险：技术债"],
         "right": ["方案 B：混合方案", "• 优势：快速上线\n• 成本：中\n• 周期：6 个月\n• 风险：依赖性"]},
        # Type D: Content — Table
        {"layout": 5, "title": "季度核心指标",
         "table": [
             {"text": "指标"}, {"text": "Q1"}, {"text": "Q2"}, {"text": "Q3"}, {"text": "Q4"},
             {"text": "营收（万）"}, {"text": "1,200"}, {"text": "1,500"}, {"text": "1,800"}, {"text": "2,100"},
             {"text": "用户（万）"}, {"text": "50"}, {"text": "68"}, {"text": "85"}, {"text": "102"},
             {"text": "NPS 得分"}, {"text": "72"}, {"text": "75"}, {"text": "78"}, {"text": "81"},
         ]},
        # Type E: Summary
        {"layout": 6, "title": "总结与展望",
         "body": ["营收同比增长 35%，超额完成目标", "核心用户突破 100 万", "技术架构升级完成",
                  "2026 年聚焦 AI 集成与海外拓展"],
         "subtitle": "感谢聆听 | 联系方式：team@example.com"},
    ]
}


# ═══════════════════════════════════════════════════════════════════════════════
#  CLI
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(
        description="PPT Toolkit v3 — parse, generate, inject, qa, demo",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__
    )
    sub = parser.add_subparsers(dest="command")

    p_parse = sub.add_parser("parse", help="Parse PPT → JSON")
    p_parse.add_argument("input", help="Input .pptx")
    p_parse.add_argument("--output", "-o", default="parsed.json")

    p_gen = sub.add_parser("generate", help="Generate PPT from JSON spec")
    p_gen.add_argument("spec", help="JSON spec file")
    p_gen.add_argument("--output", "-o", default="output.pptx")
    p_gen.add_argument("--no-badges", action="store_true", help="Skip page number badges")

    p_inj = sub.add_parser("inject", help="Template + content → PPT")
    p_inj.add_argument("template", help="Template .pptx")
    p_inj.add_argument("content", help="Content JSON")
    p_inj.add_argument("--output", "-o", default="injected.pptx")

    p_qa = sub.add_parser("qa", help="Verify PPTX for issues")
    p_qa.add_argument("input", help="Input .pptx")

    p_demo = sub.add_parser("demo", help="Generate demo (5 page types)")
    p_demo.add_argument("--output", "-o", default="demo.pptx")

    args = parser.parse_args()

    if args.command == "parse":
        data = parse_presentation(args.input)
        with open(args.output, "w", encoding="utf-8") as f:
            json.dump(data, f, indent=2, ensure_ascii=False)
        print(f"✅ Saved: {args.output}")

    elif args.command == "generate":
        with open(args.spec, "r", encoding="utf-8") as f:
            spec = json.load(f)
        prs = generate_from_spec(spec, add_badges=not args.no_badges)
        prs.save(args.output)
        print(f"✅ Saved: {args.output}")

    elif args.command == "inject":
        inject_content(args.template, args.content, args.output)
        print(f"✅ Saved: {args.output}")

    elif args.command == "qa":
        issues = qa_presentation(args.input)
        if issues:
            print(f"\n❌ Found {len(issues)} issue(s)")
            sys.exit(1)
        else:
            print(f"\n✅ QA passed (verify manually)")

    elif args.command == "demo":
        prs = generate_from_spec(DEMO_SPEC)
        prs.save(args.output)
        print(f"✅ Saved demo: {args.output}")

    else:
        parser.print_help()


if __name__ == "__main__":
    main()
